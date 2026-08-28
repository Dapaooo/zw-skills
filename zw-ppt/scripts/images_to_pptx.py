#!/usr/bin/env python3
"""Package rendered slide images into a high-fidelity, image-based PPTX."""

from __future__ import annotations

import argparse
import json
import re
from pathlib import Path

from PIL import Image
from pptx import Presentation
from pptx.util import Inches


VALID_SUFFIXES = {".png", ".jpg", ".jpeg"}


def natural_key(path: Path) -> list[object]:
    return [int(part) if part.isdigit() else part.lower() for part in re.split(r"(\d+)", path.name)]


def load_notes(path: Path | None) -> object:
    if path is None:
        return None
    return json.loads(path.read_text(encoding="utf-8"))


def note_for(notes: object, image: Path, index: int) -> list[str]:
    if notes is None:
        return []
    value = None
    if isinstance(notes, list) and index < len(notes):
        value = notes[index]
    elif isinstance(notes, dict):
        value = notes.get(image.name, notes.get(image.stem, notes.get(str(index + 1))))
    if value is None:
        return []
    if isinstance(value, list):
        return [str(item) for item in value if str(item).strip()]
    return [str(value)] if str(value).strip() else []


def set_notes(slide, paragraphs: list[str]) -> None:
    if not paragraphs:
        return
    frame = slide.notes_slide.notes_text_frame
    frame.clear()
    frame.text = paragraphs[0]
    for text in paragraphs[1:]:
        frame.add_paragraph().text = text


def main() -> None:
    parser = argparse.ArgumentParser(description=__doc__)
    parser.add_argument("image_dir", type=Path, help="Directory containing rendered slide images")
    parser.add_argument("output", type=Path, help="Output .pptx path")
    parser.add_argument("--pattern", default="slide-*", help="Image glob pattern; default: slide-*")
    parser.add_argument("--notes-json", type=Path, help="Optional JSON list or filename-to-notes mapping")
    parser.add_argument("--title", default="High-fidelity presentation")
    parser.add_argument("--allow-non-16x9", action="store_true", help="Allow images outside 16:9 and stretch them")
    args = parser.parse_args()

    images = sorted(
        [
            path
            for path in args.image_dir.glob(args.pattern)
            if path.is_file()
            and path.suffix.lower() in VALID_SUFFIXES
            and "notes" not in path.stem.lower()
        ],
        key=natural_key,
    )
    if not images:
        raise SystemExit(f"No PNG/JPG slide images found in {args.image_dir}")

    notes = load_notes(args.notes_json)
    prs = Presentation()
    prs.slide_width = Inches(13.333333)
    prs.slide_height = Inches(7.5)
    blank = prs.slide_layouts[6]
    prs.core_properties.title = args.title
    prs.core_properties.subject = "Image-based high-fidelity presentation"
    prs.core_properties.comments = "Slides are flattened page images and are not internally editable."

    for index, image_path in enumerate(images):
        with Image.open(image_path) as image:
            width, height = image.size
        ratio_error = abs((width / height) - (16 / 9))
        if ratio_error > 0.01 and not args.allow_non_16x9:
            raise SystemExit(
                f"{image_path.name} is {width}x{height}, not 16:9. "
                "Render at 1600x900 or pass --allow-non-16x9 to stretch."
            )
        slide = prs.slides.add_slide(blank)
        slide.shapes.add_picture(str(image_path), 0, 0, width=prs.slide_width, height=prs.slide_height)
        set_notes(slide, note_for(notes, image_path, index))

    args.output.parent.mkdir(parents=True, exist_ok=True)
    prs.save(args.output)
    print(f"created {args.output} with {len(images)} image-based slides")


if __name__ == "__main__":
    main()
